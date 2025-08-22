from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import io
from pptx import Presentation
import base64

app = Flask(__name__)
CORS(app)

@app.route('/status', methods=['GET'])
def status():
    return jsonify({"status": "ok"}), 200

@app.route('/extract-powerpoint', methods=['POST'])
def extract_powerpoint():
    data = request.json
    file_data = base64.b64decode(data['file_data'])
    
    try:
        prs = Presentation(io.BytesIO(file_data))
        slides_content = []
        for i, slide in enumerate(prs.slides):
            slide_content = {
                "slide_number": i + 1,
                "title": slide.shapes.title.text if slide.shapes.title else "",
                "content": [shape.text for shape in slide.shapes if shape.has_text_frame and shape is not slide.shapes.title]
            }
            slides_content.append(slide_content)
            
        return jsonify({"slides": slides_content})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/update-powerpoint', methods=['POST'])
def update_powerpoint():
    data = request.json
    original_file = data.get('original_file')
    update_instructions = data['update_instructions']
    filename = data.get('file_name', 'presentation.pptx')
    
    try:
        # Check if we're updating existing file or creating new
        if original_file and original_file != 'null' and original_file.strip():
            # Update existing presentation
            file_data = base64.b64decode(original_file)
            prs = Presentation(io.BytesIO(file_data))
            print(f"Updating existing presentation with {len(prs.slides)} slides")
        else:
            # Create new presentation from outline
            print(f"Creating new presentation from outline")
            prs = Presentation()
            
            # Remove default slide if it exists
            if len(prs.slides) > 0:
                slide_part = prs.slides._sldIdLst[0]
                prs.part.drop_rel(slide_part.rId)
                del prs.slides._sldIdLst[0]
                print("Removed default slide")
        
        # Add slides from outline
        if 'slides' in update_instructions:
            print(f"Processing {len(update_instructions['slides'])} slides from outline")
            
            for i, slide_data in enumerate(update_instructions['slides']):
                print(f"Creating slide {i+1}: {slide_data.get('title', 'No title')}")
                
                # Use available layout - check what layouts exist
                if len(prs.slide_layouts) > 1:
                    # Use title and content layout if available
                    layout = prs.slide_layouts[1]
                elif len(prs.slide_layouts) > 0:
                    # Use first available layout
                    layout = prs.slide_layouts[0]
                else:
                    # Fallback - use blank layout
                    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
                
                print(f"Using slide layout index: {prs.slide_layouts.index(layout)} out of {len(prs.slide_layouts)} available")
                slide = prs.slides.add_slide(layout)
                
                # Set title
                if 'title' in slide_data and slide.shapes.title:
                    slide.shapes.title.text = slide_data['title']
                    print(f"Set title: {slide_data['title']}")
                
                # Add content to content placeholder
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    
                    # Build content from bullets or content array
                    content_lines = []
                    if 'bullets' in slide_data and slide_data['bullets']:
                        content_lines = slide_data['bullets']
                        print(f"Added {len(content_lines)} bullet points")
                    elif 'content' in slide_data and slide_data['content']:
                        content_lines = slide_data['content']
                        print(f"Added {len(content_lines)} content lines")
                    
                    if content_lines:
                        content_placeholder.text = '\n'.join(content_lines)

        print(f"Presentation created with {len(prs.slides)} slides")
        
        # Save to memory stream
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        print(f"Sending file: {filename}")
        return send_file(
            file_stream, 
            as_attachment=True, 
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        print(f"PowerPoint generation error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(port=5001)
